import os
from pptx import Presentation
from pptx.util import Inches
from PIL import Image

# Define the folder containing your image pairs
image_folder = "your_image_folder_path"

# Initialize a PowerPoint presentation
prs = Presentation()

# Get a list of all image files in the folder
image_files = os.listdir(image_folder)

# Sort the image files to ensure that before and after images are grouped together
image_files.sort()

# Create a new slide for each pair of images
for i in range(0, len(image_files), 2):
    before_image_path = os.path.join(image_folder, image_files[i])
    after_image_path = os.path.join(image_folder, image_files[i + 1])

    # Create a slide
    slide = prs.slides.add_slide(prs.slide_layouts[5])

    # Add the "before" image
    left = Inches(0.5)
    top = Inches(1)
    width = Inches(4)
    height = Inches(3)
    slide.shapes.add_picture(before_image_path, left, top, width, height)

    # Add the "after" image
    left = Inches(5.5)
    slide.shapes.add_picture(after_image_path, left, top, width, height)

# Save the PowerPoint presentation
output_pptx = "image_pairs.pptx"
prs.save(output_pptx)

print(f"PowerPoint presentation saved as {output_pptx}")
